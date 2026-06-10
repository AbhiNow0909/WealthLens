"""
Report export — builds Excel / Word / PowerPoint artefacts from portfolio data.

The builder functions are pure: they take a `report_data` dict (+ optional
narrative) and return file bytes. Orchestration (fetching data, generating the
narrative, uploading to Supabase Storage) lives in services/export_service.py.

report_data shape:
{
  "summary": {total_value, total_invested, total_gain, total_gain_pct, portfolio_xirr},
  "holdings": [{scheme_name, isin, units_held, current_nav, invested_value, current_value, gain, gain_pct}],
  "metrics": [{isin, scheme_name, xirr, trailing_1y, alpha, beta, sharpe_ratio, sortino_ratio, max_drawdown, ...}],
  "generated_at": "YYYY-MM-DD",
}
"""

from __future__ import annotations

from io import BytesIO

from agents.state import PortfolioState


def _f(v) -> float:
    try:
        return float(v) if v is not None else 0.0
    except (TypeError, ValueError):
        return 0.0


def _pct_str(v) -> str:
    if v is None:
        return "—"
    return f"{_f(v):+.2f}%"


def _ratio_str(v) -> str:
    if v is None:
        return "—"
    return f"{_f(v):.2f}"


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def build_excel(report_data: dict) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.chart import BarChart, Reference

    s = report_data["summary"]
    holdings = report_data.get("holdings", [])
    metrics = report_data.get("metrics", [])

    header_fill = PatternFill("solid", fgColor="1A1D27")
    header_font = Font(bold=True, color="FFFFFF")
    title_font = Font(bold=True, size=14)
    money_fmt = "#,##0.00"
    pct_fmt = "+0.00;-0.00"

    wb = Workbook()

    # --- Summary sheet ---
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "WealthLens Portfolio Report"
    ws["A1"].font = title_font
    ws["A2"] = f"Generated: {report_data.get('generated_at', '')}"
    ws["A2"].font = Font(italic=True, color="666666")

    rows = [
        ("Current Value", _f(s.get("total_value"))),
        ("Invested", _f(s.get("total_invested"))),
        ("Total Gain", _f(s.get("total_gain"))),
        ("Return %", _f(s.get("total_gain_pct"))),
        ("XIRR %", _f(s.get("portfolio_xirr"))),
    ]
    r = 4
    for label, value in rows:
        ws.cell(row=r, column=1, value=label).font = Font(bold=True)
        c = ws.cell(row=r, column=2, value=value)
        c.number_format = pct_fmt if "%" in label else money_fmt
        r += 1
    ws.column_dimensions["A"].width = 22
    ws.column_dimensions["B"].width = 18

    # --- Holdings sheet ---
    wh = wb.create_sheet("Holdings")
    cols = ["Fund", "ISIN", "Units", "NAV", "Invested", "Current Value", "Gain", "Gain %"]
    for ci, name in enumerate(cols, start=1):
        cell = wh.cell(row=1, column=ci, value=name)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")
    for ri, h in enumerate(holdings, start=2):
        wh.cell(row=ri, column=1, value=h.get("scheme_name", ""))
        wh.cell(row=ri, column=2, value=h.get("isin", ""))
        wh.cell(row=ri, column=3, value=_f(h.get("units_held"))).number_format = "#,##0.000"
        wh.cell(row=ri, column=4, value=_f(h.get("current_nav"))).number_format = money_fmt
        wh.cell(row=ri, column=5, value=_f(h.get("invested_value"))).number_format = money_fmt
        wh.cell(row=ri, column=6, value=_f(h.get("current_value"))).number_format = money_fmt
        wh.cell(row=ri, column=7, value=_f(h.get("gain"))).number_format = money_fmt
        wh.cell(row=ri, column=8, value=_f(h.get("gain_pct"))).number_format = pct_fmt
    wh.column_dimensions["A"].width = 42
    wh.column_dimensions["B"].width = 16
    for col in "CDEFGH":
        wh.column_dimensions[col].width = 15

    # --- Trailing returns sheet (+ chart) ---
    wt = wb.create_sheet("Trailing Returns")
    tcols = ["Fund", "1W", "1M", "3M", "6M", "1Y", "3Y", "5Y"]
    keys = ["trailing_1w", "trailing_1m", "trailing_3m", "trailing_6m", "trailing_1y", "trailing_3y", "trailing_5y"]
    for ci, name in enumerate(tcols, start=1):
        cell = wt.cell(row=1, column=ci, value=name)
        cell.fill = header_fill
        cell.font = header_font
    for ri, m in enumerate(metrics, start=2):
        wt.cell(row=ri, column=1, value=m.get("scheme_name", ""))
        for ci, key in enumerate(keys, start=2):
            val = m.get(key)
            cell = wt.cell(row=ri, column=ci, value=_f(val) if val is not None else None)
            cell.number_format = pct_fmt
    wt.column_dimensions["A"].width = 42

    # Bar chart of 1Y returns (column F = index 6)
    if len(metrics) >= 1:
        chart = BarChart()
        chart.title = "1-Year Return by Fund"
        chart.type = "bar"
        chart.y_axis.title = "Return %"
        data_ref = Reference(wt, min_col=6, min_row=1, max_row=len(metrics) + 1)
        cats_ref = Reference(wt, min_col=1, min_row=2, max_row=len(metrics) + 1)
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        chart.height = 10
        chart.width = 24
        wt.add_chart(chart, "J2")

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------

def build_word(report_data: dict, narrative: str = "") -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor

    s = report_data["summary"]
    holdings = report_data.get("holdings", [])
    metrics = report_data.get("metrics", [])

    doc = Document()
    title = doc.add_heading("WealthLens Portfolio Report", level=0)
    title.runs[0].font.color.rgb = RGBColor(0x1A, 0x1D, 0x27)
    doc.add_paragraph(f"Generated: {report_data.get('generated_at', '')}").italic = True

    # Summary
    doc.add_heading("Summary", level=1)
    summary_tbl = doc.add_table(rows=0, cols=2)
    summary_tbl.style = "Light Grid Accent 1"
    for label, value in [
        ("Current Value", f"₹{_f(s.get('total_value')):,.2f}"),
        ("Invested", f"₹{_f(s.get('total_invested')):,.2f}"),
        ("Total Gain", f"₹{_f(s.get('total_gain')):,.2f}"),
        ("Return", _pct_str(s.get("total_gain_pct"))),
        ("XIRR", _pct_str(s.get("portfolio_xirr"))),
    ]:
        row = summary_tbl.add_row().cells
        row[0].text = label
        row[1].text = value

    # Analysis narrative
    if narrative:
        doc.add_heading("Analysis", level=1)
        for para in narrative.split("\n"):
            para = para.strip()
            if para:
                doc.add_paragraph(para)

    # Holdings
    doc.add_heading("Holdings", level=1)
    htbl = doc.add_table(rows=1, cols=5)
    htbl.style = "Light Grid Accent 1"
    for i, name in enumerate(["Fund", "Invested", "Current Value", "Gain", "Gain %"]):
        htbl.rows[0].cells[i].text = name
    for h in holdings:
        c = htbl.add_row().cells
        c[0].text = h.get("scheme_name", "")
        c[1].text = f"₹{_f(h.get('invested_value')):,.0f}"
        c[2].text = f"₹{_f(h.get('current_value')):,.0f}"
        c[3].text = f"₹{_f(h.get('gain')):,.0f}"
        c[4].text = _pct_str(h.get("gain_pct"))

    # Risk metrics
    doc.add_heading("Risk Metrics", level=1)
    mtbl = doc.add_table(rows=1, cols=6)
    mtbl.style = "Light Grid Accent 1"
    for i, name in enumerate(["Fund", "Alpha", "Beta", "Sharpe", "Sortino", "Max DD"]):
        mtbl.rows[0].cells[i].text = name
    for m in metrics:
        c = mtbl.add_row().cells
        c[0].text = m.get("scheme_name", "")
        c[1].text = _pct_str(m.get("alpha"))
        c[2].text = _ratio_str(m.get("beta"))
        c[3].text = _ratio_str(m.get("sharpe_ratio"))
        c[4].text = _ratio_str(m.get("sortino_ratio"))
        c[5].text = _pct_str(m.get("max_drawdown"))

    # Shrink table fonts a touch for density
    for tbl in (htbl, mtbl):
        for row in tbl.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def build_ppt(report_data: dict, narrative: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    s = report_data["summary"]
    holdings = report_data.get("holdings", [])
    metrics_by_isin = {m.get("isin"): m for m in report_data.get("metrics", [])}

    BG = RGBColor(0x0F, 0x11, 0x17)
    SURFACE = RGBColor(0x1A, 0x1D, 0x27)
    TEXT = RGBColor(0xF1, 0xF5, 0xF9)
    MUTED = RGBColor(0x94, 0xA3, 0xB8)
    ACCENT = RGBColor(0x4F, 0x8E, 0xF7)
    GAIN = RGBColor(0x22, 0xC5, 0x5E)
    LOSS = RGBColor(0xEF, 0x44, 0x44)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

    def add_text(slide, text, left, top, width, height, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    # --- Title slide ---
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_text(slide, "WealthLens", 0.8, 2.4, 11.7, 1.0, size=44, color=ACCENT, bold=True)
    add_text(slide, "Portfolio Report", 0.8, 3.4, 11.7, 0.8, size=28, color=TEXT)
    add_text(slide, f"Generated {report_data.get('generated_at', '')}", 0.8, 4.2, 11.7, 0.5, size=14, color=MUTED)

    # --- Summary slide ---
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_text(slide, "Portfolio Summary", 0.8, 0.6, 11.7, 0.8, size=28, color=TEXT, bold=True)
    gain = _f(s.get("total_gain"))
    cards = [
        ("Current Value", f"₹{_f(s.get('total_value')):,.0f}", TEXT),
        ("Invested", f"₹{_f(s.get('total_invested')):,.0f}", TEXT),
        ("Total Gain", f"₹{gain:,.0f}", GAIN if gain >= 0 else LOSS),
        ("Return", _pct_str(s.get("total_gain_pct")), GAIN if _f(s.get("total_gain_pct")) >= 0 else LOSS),
    ]
    for i, (label, value, color) in enumerate(cards):
        left = 0.8 + i * 3.0
        card = slide.shapes.add_shape(1, Inches(left), Inches(2.0), Inches(2.7), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.fill.background()
        add_text(slide, label, left + 0.2, 2.2, 2.4, 0.5, size=13, color=MUTED)
        add_text(slide, value, left + 0.2, 2.8, 2.4, 0.8, size=22, color=color, bold=True)

    # --- One slide per fund ---
    for h in holdings:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        m = metrics_by_isin.get(h.get("isin"), {})
        add_text(slide, h.get("scheme_name", ""), 0.8, 0.6, 11.7, 1.0, size=22, color=TEXT, bold=True)
        add_text(slide, h.get("isin", ""), 0.8, 1.5, 11.7, 0.4, size=12, color=MUTED)
        gpct = _f(h.get("gain_pct"))
        lines = [
            ("Current Value", f"₹{_f(h.get('current_value')):,.0f}", TEXT),
            ("Invested", f"₹{_f(h.get('invested_value')):,.0f}", TEXT),
            ("Gain", _pct_str(h.get("gain_pct")), GAIN if gpct >= 0 else LOSS),
            ("1Y Return", _pct_str(m.get("trailing_1y")), TEXT),
            ("Sharpe", _ratio_str(m.get("sharpe_ratio")), TEXT),
            ("Max Drawdown", _pct_str(m.get("max_drawdown")), LOSS),
        ]
        for i, (label, value, color) in enumerate(lines):
            row = i // 3
            col = i % 3
            left = 0.8 + col * 4.0
            top = 2.4 + row * 1.6
            add_text(slide, label, left, top, 3.6, 0.4, size=13, color=MUTED)
            add_text(slide, value, left, top + 0.45, 3.6, 0.7, size=22, color=color, bold=True)

    # --- Recommendation slide ---
    if narrative:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_text(slide, "Analysis & Recommendation", 0.8, 0.6, 11.7, 0.8, size=26, color=ACCENT, bold=True)
        text = narrative.strip()
        if len(text) > 1400:
            text = text[:1400].rsplit(" ", 1)[0] + "…"
        add_text(slide, text, 0.8, 1.6, 11.7, 5.4, size=14, color=TEXT)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# LangGraph node
# ---------------------------------------------------------------------------

async def run(state: PortfolioState) -> PortfolioState:
    """
    Generate a report when the graph routes here. Format is taken from
    state["export_format"] (default 'excel'); the heavy lifting is delegated
    to the export service.
    """
    from services.export_service import generate_report

    fmt = state.get("export_format", "excel") if isinstance(state, dict) else "excel"
    try:
        result = await generate_report(state["user_id"], fmt, state.get("user_prompt", ""))
        return {**state, "export_path": result.get("export_url")}
    except Exception:
        return {**state, "export_path": None}
